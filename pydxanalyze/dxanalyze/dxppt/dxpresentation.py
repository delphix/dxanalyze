#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
#
# Copyright (c) 2019 by Delphix. All rights reserved.
#

import time
import tempfile
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

import dxanalyze.dxppt.dxslideconfig as dxslideconfig


def pptx_delete_slide(prs, slide):
    """
    pptx package doesn't provide a delete slide function
    workaround from github issue - https://github.com/scanny/python-pptx/issues/67#issuecomment-296135015
    :param1 prs: Presentaton object
    :param2 slide: Slide object to delete
    """

    # Make dictionary with necessary information
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

def delete_slides(prs, delete_list):
    """
    remove slides from presentation based on list of analytics missing 
    maninly iscsi or nfs will be deleted but it support any analytics
    :param1 prs: Presentaton object
    :param2 delete_list: List of analytics name to delete 
    """
    slide_no_delete_list = []
    for analytic_name in delete_list:
        slide_no_delete_list.extend(list(dxslideconfig.slide_with_pictures[analytic_name]))
    # no of slides to delete has to be ordered in descending order
    # as removing a slide before is changing a numbers
    slide_no_delete_list.sort(reverse = True)
    for slide_no in slide_no_delete_list:
        pptx_delete_slide(prs, prs.slides[slide_no-1])


def update_titles(prs, dlpx_engine_name, author_name):
    """
    Update title on each slide plus update a first slide
    :param1 prs: Presentaton object
    :param2 dlpx_engine_name: Engine name 
    :param3 author_name: name of enginner to show on 1st page
    """
    for slide in prs.slides: # iterate over each slide
        # title_shape =  slide.shapes[0] # consider the zeroth indexed shape as the title
        # if title_shape.has_text_frame: # is this shape has textframe attribute true then
        slidenum = prs.slides.index(slide) + 1
        if slidenum == 1:
            slideNameDate = slide.shapes[0]
            slideHeading = slide.shapes[1]
            slideHeading.text_frame.paragraphs[0].runs[0].text = dlpx_engine_name
            if slideNameDate.has_text_frame:
                runs = slideNameDate.text_frame.paragraphs[0].runs
                runs[0].text = author_name
                runs[0].font.color.rgb = RGBColor(26, 214, 245)
                runs[7].text = time.strftime("%d %b %Y")
                runs[7].font.color.rgb = RGBColor(255, 255, 255)
        elif slidenum not in dxslideconfig.skip_title_update_slides:
            if slide.shapes.title:
                slide.shapes.title.text = dlpx_engine_name + " " + slide.shapes.title.text

def add_pictures(prs, analytic_list, imgdir):
    """
    Add generated picture from directory into slide
    :param1 prs: Presentaton object
    :param2 analytic_list: List of analytics where pictures will be added 
    :param3 imgdir: picture directory
    """
    for analytic_name in analytic_list:
        analytic_graphs = dxslideconfig.slide_with_pictures[analytic_name]
        for slide_no, graph_name in analytic_graphs.items():
            fname = os.path.join(imgdir, graph_name)
            if os.path.isfile(fname):                
                slide = prs.slides[slide_no-1]
                left = Inches(0.2)
                top = Inches(1.1)
                width = Inches(6.5)
                height = Inches(4.0)
                slide.shapes.add_picture(fname, left, top, width, height)

def cleanup_pictures(analytic_list, imgdir):
    """
    Delete generated pictures
    :param1 analytic_list: List of analytics to list pictures to deletion
    :param2 imgdir: picture directory
    """
    for analytic_name in analytic_list:
        analytic_graphs = dxslideconfig.slide_with_pictures[analytic_name]
        for slide_no, graph_name in analytic_graphs.items():
            fname = os.path.join(imgdir, graph_name)
            #print("testing {}".format(fname))
            if os.path.isfile(fname):
                #print("deleting {}".format(fname))
                os.remove(fname)

def gen_presentation(analytic_list, out_location, engine_name):
    """
    Generate presentation based on the template and save it as a new one
    :param1 analytic_list: List of analytics with data to add to presentation
    :param2 out_location: output directory to save presentation
    :param3 engine_name: Delphix Engine name 
    """
    prs = Presentation(dxslideconfig.report_template)
    update_titles(prs, engine_name, "")
    # a cache hit ration needs to be added to a list to include it in end report
    analytic_list.append("chr")
    tempdir = tempfile.gettempdir()
    add_pictures(prs, analytic_list, tempdir)

    delete_slide_list = []

    for analytic in ["cpu", "network", "iscsi", "nfs", "disk"]:
        # for analytics which are empty, we need to delete slides
        if analytic not in analytic_list:
            delete_slide_list.append(analytic)
            
    delete_slides(prs, delete_slide_list)

    fname = os.path.join(out_location, "{}_analytics.pptx".format(engine_name))
    prs.save(fname)

    #cleanup_pictures(analytic_list, tempdir)
    print ("Report {} generated.".format(fname))

def gen_farm_presentation(out_location):
    """
    Generate presentation based on the template and save it as a new one
    :param1 analytic_list: List of analytics with data to add to presentation
    :param2 out_location: output directory to save presentation
    :param3 engine_name: Delphix Engine name 
    """
    prs = Presentation(dxslideconfig.farm_report_template)
    #update_titles(prs, "Farm Engine", "Ajay T")
    tempdir = tempfile.gettempdir()
    analytic_list = ['farm']
    add_pictures(prs, analytic_list, tempdir)

    fname = os.path.join(out_location, "{}_analytics.pptx".format("farm"))
    prs.save(fname)

    #cleanup_pictures(analytic_list, tempdir)
    print ("Report {} generated.".format(fname))
